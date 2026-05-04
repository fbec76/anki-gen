from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation


def extract_text(file_path: Path) -> str:
    """Dispatch to the right extractor based on file extension."""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext in {".pptx", ".ppt"}:
        return _extract_pptx(file_path)
    if ext in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8")
    raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        parts.append(f"\n--- Slide/Page {i} ---\n{text}")
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        buf = [f"\n--- Slide {i} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                buf.append(shape.text)
        parts.append("\n".join(buf))
    return "\n".join(parts)